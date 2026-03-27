from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
DOCS = ROOT / "docs"
OUTPUT = DOCS / "SmartCourier-Project-Presentation.pptx"
ARCHITECTURE_SVG = DOCS / "architecture-diagram.svg"

TITLE_COLOR = RGBColor(15, 56, 103)
ACCENT_COLOR = RGBColor(34, 139, 230)
TEXT_COLOR = RGBColor(40, 40, 40)
MUTED_COLOR = RGBColor(98, 110, 125)
LIGHT_BG = RGBColor(245, 248, 252)


def set_bg(slide, color=RGBColor(255, 255, 255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.05), Inches(11.6), Inches(0.45))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Aptos"
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED_COLOR

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.45), Inches(2.1), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()


def add_bullets(slide, items, left=0.9, top=1.7, width=5.8, height=4.8, font_size=19):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)


def add_two_column(slide, left_title, left_items, right_title, right_items):
    left_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.8), Inches(5.8), Inches(4.8)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = LIGHT_BG
    left_card.line.color.rgb = RGBColor(220, 228, 238)

    right_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = LIGHT_BG
    right_card.line.color.rgb = RGBColor(220, 228, 238)

    lt = slide.shapes.add_textbox(Inches(1.0), Inches(2.05), Inches(5.1), Inches(0.4))
    lp = lt.text_frame.paragraphs[0]
    lr = lp.add_run()
    lr.text = left_title
    lr.font.name = "Aptos Display"
    lr.font.size = Pt(20)
    lr.font.bold = True
    lr.font.color.rgb = TITLE_COLOR

    rt = slide.shapes.add_textbox(Inches(7.1), Inches(2.05), Inches(5.1), Inches(0.4))
    rp = rt.text_frame.paragraphs[0]
    rr = rp.add_run()
    rr.text = right_title
    rr.font.name = "Aptos Display"
    rr.font.size = Pt(20)
    rr.font.bold = True
    rr.font.color.rgb = TITLE_COLOR

    add_bullets(slide, left_items, left=1.0, top=2.55, width=5.0, height=3.7, font_size=16)
    add_bullets(slide, right_items, left=7.1, top=2.55, width=5.0, height=3.7, font_size=16)


def add_service_grid(slide, services):
    positions = [
        (0.7, 1.8), (4.4, 1.8), (8.1, 1.8),
        (0.7, 4.2), (4.4, 4.2), (8.1, 4.2),
    ]
    for (title, body), (x, y) in zip(services, positions):
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.2), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(220, 228, 238)

        tbox = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.15), Inches(2.8), Inches(0.4))
        p = tbox.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = TITLE_COLOR

        bbox = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.55), Inches(2.8), Inches(1.0))
        bp = bbox.text_frame.paragraphs[0]
        br = bp.add_run()
        br.text = body
        br.font.name = "Aptos"
        br.font.size = Pt(12)
        br.font.color.rgb = TEXT_COLOR


def add_architecture_slide(slide):
    add_title(slide, "System Architecture", "Platform services, business services, databases, messaging, and tracing")
    if ARCHITECTURE_SVG.exists():
        try:
            slide.shapes.add_picture(str(ARCHITECTURE_SVG), Inches(0.55), Inches(1.65), width=Inches(12.2))
            return
        except Exception:
            pass

    add_two_column(
        slide,
        "Platform Layer",
        [
            "API Gateway routes all external traffic to internal services.",
            "Eureka registers and discovers service instances dynamically.",
            "Config Server serves centralized YAML configuration from config-repo.",
            "Zipkin receives distributed traces from gateway and all services.",
        ],
        "Business Layer",
        [
            "Auth, Delivery, Tracking, and Admin are isolated Spring Boot services.",
            "Each service owns its own MySQL database for data autonomy.",
            "RabbitMQ handles asynchronous delivery events across services.",
            "OpenFeign + LoadBalancer power synchronous admin-to-delivery calls.",
        ],
    )


def add_demo_flow(slide):
    add_title(slide, "Suggested Live Demo Flow", "A short sequence that shows the whole platform clearly")
    steps = [
        "1. Open Eureka dashboard and show registered services.",
        "2. Open gateway Swagger and show the service dropdown.",
        "3. Sign up or log in as customer through auth-service.",
        "4. Log in as admin and keep the JWT token ready.",
        "5. Create a delivery through delivery-service.",
        "6. Fetch admin delivery overview to show Feign-based synchronous communication.",
        "7. Update delivery status to IN_TRANSIT or DELAYED.",
        "8. Show tracking updates and exception handling created from RabbitMQ events.",
        "9. Open Zipkin and show the distributed trace.",
        "10. Close with JaCoCo, SonarQube, and Docker-based deployment.",
    ]
    add_bullets(slide, steps, left=0.9, top=1.8, width=11.4, height=5.2, font_size=18)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_BG)
    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = TITLE_COLOR
    banner.line.fill.background()
    t = slide.shapes.add_textbox(Inches(0.65), Inches(1.55), Inches(10.5), Inches(1.4))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "SmartCourier Delivery Management System"
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT
    s = t.text_frame.add_paragraph()
    s.text = (
        "Spring Boot Microservices Project Presentation\n"
        "API Gateway, Eureka, Config Server, JWT, OpenFeign, RabbitMQ, Zipkin, Docker, JaCoCo, and SonarQube"
    )
    s.font.name = "Aptos"
    s.font.size = Pt(17)
    s.font.color.rgb = MUTED_COLOR
    info = slide.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(11.6), Inches(1.2))
    ip = info.text_frame.paragraphs[0]
    ir = ip.add_run()
    ir.text = (
        "Prepared for project viva and architecture walkthrough.\n"
        "This presentation explains the system from problem statement to deployment, testing, and tracing."
    )
    ir.font.name = "Aptos"
    ir.font.size = Pt(20)
    ir.font.color.rgb = TEXT_COLOR

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Project Overview", "What the system does and why a microservices approach was used")
    add_two_column(
        slide,
        "Business Goal",
        [
            "Handle user signup, login, courier booking, tracking, status updates, and admin operations.",
            "Support delivery exceptions, proof of delivery, reporting, and operational monitoring.",
            "Provide a realistic logistics workflow instead of isolated CRUD services.",
        ],
        "Why Microservices",
        [
            "Separate authentication, delivery, tracking, and admin responsibilities.",
            "Improve maintainability, scalability, and fault isolation.",
            "Allow both synchronous and asynchronous interservice communication patterns.",
        ],
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_architecture_slide(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Technology Stack", "Core frameworks and infrastructure used in the project")
    add_two_column(
        slide,
        "Application Stack",
        [
            "Spring Boot 3 and Spring Cloud for microservices.",
            "Spring Security + JWT for authentication and authorization.",
            "Spring Data JPA + Hibernate for persistence.",
            "Springdoc OpenAPI / Swagger UI for API documentation.",
            "Resilience4j Circuit Breaker and Caffeine Cache in admin flows.",
        ],
        "Infrastructure Stack",
        [
            "API Gateway for central routing.",
            "Eureka for service discovery.",
            "Config Server + config-repo for centralized configuration.",
            "RabbitMQ for asynchronous event-driven communication.",
            "Zipkin for distributed tracing, Docker Compose for deployment, SonarQube + JaCoCo for quality.",
        ],
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Service Responsibilities", "Each service has a focused responsibility and its own database")
    add_service_grid(
        slide,
        [
            ("API Gateway", "Single entry point, route forwarding, centralized Swagger access, and discovery-based routing."),
            ("Eureka Server", "Registers services and helps gateway and Feign discover live instances dynamically."),
            ("Config Server", "Serves centralized YAML configuration from config-repo to all services."),
            ("Auth Service", "User signup, login, JWT generation, role-based access, and seeded admin account."),
            ("Delivery Service", "Delivery creation, charge calculation, tracking number generation, and status lifecycle."),
            ("Tracking + Admin", "Tracking timeline, proof/documents, dashboard, reports, exceptions, and operational overview."),
        ],
    )

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Interservice Communication", "The project uses both synchronous and asynchronous integration styles")
    add_two_column(
        slide,
        "Synchronous Flow",
        [
            "Admin-service uses OpenFeign to call delivery-service.",
            "Eureka + Spring Cloud LoadBalancer resolve the logical service name.",
            "Circuit breaker protects admin flows if delivery-service is unavailable.",
            "Used for live overviews, dashboard data, and direct operational reads.",
        ],
        "Asynchronous Flow",
        [
            "Delivery-service publishes lifecycle events to RabbitMQ.",
            "Tracking-service consumes events to create tracking history.",
            "Admin-service consumes delayed/failed/returned events for exception handling.",
            "This reduces coupling and supports eventual consistency across services.",
        ],
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Security and Access Control", "Authentication is centralized while authorization is enforced service by service")
    add_bullets(
        slide,
        [
            "Auth-service manages signup, login, password encoding, JWT issuance, and default admin creation.",
            "Clients send JWT tokens in the Authorization header using the Bearer scheme.",
            "Each service validates JWTs and enforces role-based authorization such as ROLE_CUSTOMER and ROLE_ADMIN.",
            "SecurityConfig classes define protected routes, public endpoints, and CORS handling.",
            "Gateway makes testing easier, but security remains enforced inside each business service.",
        ],
        left=0.85,
        top=1.8,
        width=11.5,
        height=4.9,
        font_size=19,
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Data Management and Caching", "Each service owns its data while admin reads are optimized")
    add_two_column(
        slide,
        "Database Design",
        [
            "Every business service uses its own MySQL database.",
            "This avoids tight coupling and follows the database-per-service microservices pattern.",
            "JPA entities, repositories, and service layers keep persistence clean and maintainable.",
        ],
        "Caching",
        [
            "Admin-service uses Caffeine for repeated reads such as dashboard and overview responses.",
            "Cache reduces unnecessary repeated work and improves response time for read-heavy operations.",
            "Write/update operations can invalidate or refresh cached data when required.",
        ],
    )

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "API Gateway and Unified Swagger", "All service APIs can now be accessed from one gateway UI")
    add_bullets(
        slide,
        [
            "Gateway Swagger UI shows a dropdown for gateway, auth-service, delivery-service, tracking-service, and admin-service.",
            "OpenAPI docs for downstream services are proxied through the gateway using /v3/api-docs/... routes.",
            "Native /api/v1/... pass-through routes were added so 'Try it out' works without opening each service port manually.",
            "This makes the system easier to demo, test, and explain during viva.",
        ],
        left=0.85,
        top=1.8,
        width=11.6,
        height=4.6,
        font_size=19,
    )

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Distributed Tracing with Zipkin", "Tracing helps visualize one request across gateway and downstream services")
    add_two_column(
        slide,
        "How It Works",
        [
            "Micrometer tracing creates trace IDs and span IDs for requests.",
            "Services export spans to Zipkin using the configured endpoint.",
            "Gateway, Config Server, Eureka, and all business services contribute trace data.",
        ],
        "Why It Matters",
        [
            "Shows request journey across microservices in one timeline.",
            "Helps debug latency, service dependencies, and failures.",
            "Useful during live demos to prove distributed system behavior.",
        ],
    )

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Deployment and Run Workflow", "The project is designed for a Docker-first demonstration flow")
    add_bullets(
        slide,
        [
            "Full rebuild after major changes: ./scripts/run-all-services.sh --build",
            "Normal startup when images already exist: ./scripts/run-all-services.sh",
            "Selective rebuild for changed services: ./scripts/rebuild-services.sh <service-names>",
            "Stop all containers: ./scripts/stop-all-services.sh",
            "Restart the platform: ./scripts/restart-all-services.sh",
            "Core runtime URLs: gateway 8080, Eureka 8761, Config Server 8888, RabbitMQ 15672, Zipkin 9411, SonarQube 9000",
        ],
        left=0.8,
        top=1.75,
        width=11.7,
        height=5.1,
        font_size=18,
    )

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Testing and Quality Assurance", "The project includes automated testing, coverage, and static analysis")
    add_two_column(
        slide,
        "Testing and Coverage",
        [
            "JUnit and Mockito are used for service-layer and controller-oriented test cases.",
            "mvn test verifies whether tests pass; BUILD SUCCESS confirms success.",
            "mvn -Pcoverage verify generates JaCoCo reports for coverage analysis.",
        ],
        "Quality Analysis",
        [
            "SonarQube reads the codebase and JaCoCo XML reports.",
            "It highlights bugs, vulnerabilities, code smells, duplication, and maintainability.",
            "Together, JaCoCo and SonarQube show both test depth and code quality.",
        ],
    )

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_demo_flow(slide)

    # Slide 14
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Key Strengths and Conclusion", "Why this project stands out as a complete backend microservices platform")
    add_bullets(
        slide,
        [
            "Clear separation of concerns across auth, delivery, tracking, and admin services.",
            "Real enterprise features: API Gateway, Eureka, Config Server, Feign, RabbitMQ, Circuit Breaker, Caching, Zipkin, Docker.",
            "Security with JWT and role-based access control.",
            "Operational quality with tests, JaCoCo coverage, and SonarQube analysis.",
            "Ready for both GitHub portfolio presentation and academic viva demonstration.",
        ],
        left=0.85,
        top=1.85,
        width=11.4,
        height=4.8,
        font_size=20,
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    create_presentation()
